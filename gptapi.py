import PyPDF2
import os
import requests
from dotenv import load_dotenv

import re
from pptx import Presentation
import tiktoken

# Load necessary components
  # Ensure 'punkt' data is downloaded
load_dotenv()  # Load environment variables from .env file
encoder = tiktoken.encoding_for_model("gpt-4o-mini")
# Get OpenAI API Key from .env file
OPENAI_API_KEY = os.getenv('OPEN_API_KEY')
prompts = {}
prompts['exam_sum'] = "Summarize the following text in order to prepare for an exam, break it down to sub topics and put bullet points under each topic, make sure every topic is covered, while giving output make sure every Bold sentnce is in the format **sentence** only make Headings and sub headings in bold"
prompts['qp_dup'] = "Make another question paper following the same format and the same level of difficulty for the same subject, keep the ratio of theory to numerical/coding questions the same "
prompts['question'] = "From the following text generate various relevant questions to test the understanding of the student generate 5 MCQ questions, 5 2 mark questions 5 3 mark questions and 5 5 mark questions "
# Function to split text into manageable chunks


def count_tokens(text):
    return len(encoder.encode(text))
def split_text_into_chunks(text, max_tokens=3000):
    sentences = re.split(r'(?<!\w\.\w.)(?<![A-Z][a-z]\.)(?<=\.|\?)\s', text)
    chunks = []
    current_chunk = []
    current_tokens = 0

    for sentence in sentences:
        sentence_tokens = count_tokens(sentence)
        if current_tokens + sentence_tokens <= max_tokens:
            current_chunk.append(sentence)
            current_tokens += sentence_tokens
        else:
            chunks.append(" ".join(current_chunk))
            current_chunk = [sentence]
            current_tokens = sentence_tokens

    if current_chunk:
        chunks.append(" ".join(current_chunk))

    return chunks


# Function to summarize text using OpenAI API endpoint
def summarize_text_with_api(chunk, max_tokens=3000):
    api_url = "https://api.openai.com/v1/chat/completions"
    headers = {
        "Authorization": f"Bearer {OPENAI_API_KEY}",
        "Content-Type": "application/json"
    }

    payload = {
        "model": "gpt-4o-mini",
        "messages": [
            {"role": "system", "content": "You are an expert summarizer."},
            {"role": "user", "content": f"{prompts['exam_sum']}:\n\n{chunk}"}
        ],
        "max_tokens": max_tokens,
        "temperature": 0.7
    }

    response = requests.post(api_url, headers=headers, json=payload)

    if response.status_code == 200:
        return response.json()['choices'][0]['message']['content'].strip(),count_tokens(chunk)
    else:
        print(f"Error: {response.status_code} - {response.text}")
        return "[Error occurred during summarization]"

def generate_questions_with_api(chunk, max_tokens=3000):
    api_url = "https://api.openai.com/v1/chat/completions"
    headers = {
        "Authorization": f"Bearer {OPENAI_API_KEY}",
        "Content-Type": "application/json"
    }

    payload = {
        "model": "gpt-4o-mini",
        "messages": [
            {"role": "system", "content": "You are an expert question setter for a university exam."},
            {"role": "user", "content": f"{prompts['question']}:\n\n{chunk}"}
        ],
        "max_tokens": max_tokens,
        "temperature": 0.7
    }

    response = requests.post(api_url, headers=headers, json=payload)

    if response.status_code == 200:
        return response.json()['choices'][0]['message']['content'].strip(),count_tokens(chunk)
    else:
        print(f"Error: {response.status_code} - {response.text}")
        return "[Error occurred during summarization]"


def gen_ques_from_pdf(pdf_path, Save_to_txt=False):
    extracted_text = ""
    total_tokens_used = 0
    try:
        with open(pdf_path, 'rb') as pdf_file:
            reader = PyPDF2.PdfReader(pdf_file)
            for page in reader.pages:
                text = page.extract_text()
                extracted_text += text.strip() if text else "[No text found on this page]\n"
        print(extracted_text)
    except Exception as e:
        print(f"An error occurred while reading the PDF: {e}")
        return None

    # Split text into manageable chunks
    chunks = split_text_into_chunks(extracted_text, max_tokens=3000)

    # Summarize each chunk using OpenAI API
    questions = ""
    for idx, chunk in enumerate(chunks):
        print(f"Generating Questions chunk {idx + 1}/{len(chunks)}...")
        chunk_questions, chunk_tokens = generate_questions_with_api(chunk)
        questions += chunk_questions + "\n"
        total_tokens_used += chunk_tokens

    # Optionally save the summary to a text file


    return questions,total_tokens_used

def summarize_pdf(pdf_path, Save_to_txt=False):
    extracted_text = ""
    total_tokens_used = 0

    try:
        with open(pdf_path, 'rb') as pdf_file:
            reader = PyPDF2.PdfReader(pdf_file)
            for page in reader.pages:
                text = page.extract_text()
                extracted_text += text.strip() if text else "[No text found on this page]\n"
    except Exception as e:
        print(f"An error occurred while reading the PDF: {e}")
        return None, total_tokens_used

    # Split text into manageable chunks
    chunks = split_text_into_chunks(extracted_text, max_tokens=3000)

    # Summarize each chunk
    summary = ""
    for idx, chunk in enumerate(chunks):
        print(f"Summarizing chunk {idx + 1}/{len(chunks)}...")
        chunk_summary, chunk_tokens = summarize_text_with_api(chunk)
        summary += chunk_summary + "\n"
        total_tokens_used += chunk_tokens

    # Optionally save the summary to a text file
    if Save_to_txt:
        try:
            output_file = f"{os.path.splitext(pdf_path)[0]}_summary.txt"
            with open(output_file, 'w') as file:
                file.write(summary)
            print(f"Summary has been saved to {output_file}")
        except Exception as e:
            print(f"An error occurred while saving the summary: {e}")

    return summary, total_tokens_used

def summarize_ppt(ppt_path, Save_to_txt=False):
    extracted_text = ""
    total_tokens_used = 0

    try:
        presentation = Presentation(ppt_path)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    extracted_text += shape.text.strip() + "\n"
    except Exception as e:
        print(f"An error occurred while reading the PowerPoint: {e}")
        return None, total_tokens_used

    # Split text into manageable chunks
    chunks = split_text_into_chunks(extracted_text, max_tokens=3000)

    # Summarize each chunk
    summary = ""
    for idx, chunk in enumerate(chunks):
        print(f"Summarizing chunk {idx + 1}/{len(chunks)}...")
        chunk_summary, chunk_tokens = summarize_text_with_api(chunk)
        summary += chunk_summary + "\n"
        total_tokens_used += chunk_tokens

    # Optionally save the summary to a text file
    if Save_to_txt:
        try:
            output_file = f"{os.path.splitext(ppt_path)[0]}_summary.txt"
            with open(output_file, 'w') as file:
                file.write(summary)
            print(f"Summary has been saved to {output_file}")
        except Exception as e:
            print(f"An error occurred while saving the summary: {e}")

    return summary, total_tokens_used

def gen_ques_from_ppt(ppt_path, Save_to_txt=False):
    """Extract text from a PowerPoint and generate questions."""
    extracted_text = ""
    total_tokens_used=0

    try:
        presentation = Presentation(ppt_path)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    extracted_text += shape.text.strip() + "\n"
    except Exception as e:
        print(f"An error occurred while reading the PowerPoint: {e}")
        return None

    # Split text into manageable chunks
    chunks = split_text_into_chunks(extracted_text, max_tokens=3000)

    # Generate questions for each chunk using the API
    questions = ""
    for idx, chunk in enumerate(chunks):
        print(f"Generating questions for chunk {idx + 1}/{len(chunks)}...")
        chunk_questions,chunk_tokens= generate_questions_with_api(chunk)
        questions += chunk_questions + "\n"
        total_tokens_used += chunk_tokens

    # Optionally save the questions to a text file
    if Save_to_txt:
        try:
            output_file = f"{os.path.splitext(ppt_path)[0]}_questions.txt"
            with open(output_file, 'w') as file:
                file.write(questions)
            print(f"Questions have been saved to {output_file}")
        except Exception as e:
            print(f"An error occurred while saving the questions: {e}")

    return questions,total_tokens_used

#pdf_path = "example.pdf"
#summary = summarize_pdf(pdf_path, Save_to_txt=True)
#print("Summary:")
#print(summary)
