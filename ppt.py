from pptx import Presentation
from dotenv import load_dotenv
import os
import requests
import nltk
from nltk.tokenize import sent_tokenize
from transformers import pipeline, BartTokenizer

# Load necessary components
nltk.download('punkt')  # Ensure 'punkt' data is downloaded
load_dotenv()  # Load environment variables from .env file

# Get Hugging Face API Key from .env file
HF_API_KEY = os.getenv('HUGGING_API_KEY')
print(HF_API_KEY)

# Define the API URL for Hugging Face's Inference API
API_URL = "https://api-inference.huggingface.co/models/facebook/bart-large-cnn"

# Initialize Hugging Face pipeline for summarization (local)
summarizer = pipeline("summarization", model="facebook/bart-large-cnn")
tokenizer = BartTokenizer.from_pretrained("facebook/bart-large-cnn")


# Function to split text efficiently
def split_text_into_chunks(text, max_tokens=1024):
    sentences = sent_tokenize(text)
    chunks = []
    current_chunk = []
    current_tokens = 0

    for sentence in sentences:
        tokenized_sentence = tokenizer(sentence, truncation=False, return_tensors="pt")["input_ids"]
        sentence_length = tokenized_sentence.shape[1]  # Number of tokens in the sentence

        if current_tokens + sentence_length <= max_tokens:
            current_chunk.append(sentence)
            current_tokens += sentence_length
        else:
            chunks.append(" ".join(current_chunk))
            current_chunk = [sentence]
            current_tokens = sentence_length

    if current_chunk:
        chunks.append(" ".join(current_chunk))

    return chunks


# Function to summarize text using Hugging Face Inference API (remote)
def summarize_text_with_api(text, max_length=500, min_length=300):
    headers = {
        "Authorization": f"Bearer {HF_API_KEY}"
    }

    payload = {
        "inputs": text,
        "parameters": {
            "max_length": max_length,
            "min_length": min_length,
            "do_sample": False
        }
    }

    response = requests.post(API_URL, headers=headers, json=payload)

    print(f"Response Status Code: {response.status_code}")
    print(f"Response Content: {response.text}")

    if response.status_code == 200:
        summary = response.json()
        return summary[0]['summary_text']
    else:
        print(f"Error: {response.status_code}")
        return f"Error occurred: {response.text}"


# Function to summarize text locally using Hugging Face model (pipeline)
def summarize_text_locally(text, max_length=750, min_length=500):
    summary = summarizer(text, max_length=max_length, min_length=min_length, do_sample=False)
    return summary[0]['summary_text']


# Main function to extract, split, and summarize PowerPoint files
def summarize_ppt(ppt_path, USE_API=False, Save_to_txt=False):
    extracted_text = ""

    try:
        # Open the PowerPoint file and extract text
        ppt = Presentation(ppt_path)
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    extracted_text += shape.text.strip() + "\n"

    except Exception as e:
        print(f"An error occurred while extracting text: {e}")
        return None

    # Tokenize the extracted text and check if it exceeds the token limit
    print(extracted_text)
    tokens = tokenizer(extracted_text, truncation=False)["input_ids"]
    print(f"Number of tokens: {len(tokens)}")

    if len(tokens) > 1024:
        print("The text exceeds the maximum token limit and needs to be split.")

    # Split text if token count exceeds the limit
    chunks = split_text_into_chunks(extracted_text, max_tokens=1024)

    # Generate summaries for each chunk
    summary = ""
    for idx, chunk in enumerate(chunks):
        print(f"Summarizing chunk {idx + 1}...")

        if USE_API:
            summary += summarize_text_with_api(chunk) + "\n" + "\n"
        else:
            summary += summarize_text_locally(chunk) + "\n" + "\n"

    # Save the summary to a text file if required
    if Save_to_txt:
        try:
            with open(f"{ppt_path[:-5]}_summary.txt", 'w', encoding='utf-8') as file:
                file.write(summary)
            print(f"Summary has been saved to {ppt_path[:-5]}_summary.txt")
        except Exception as e:
            print(f"An error occurred while saving the summary: {e}")

    return summary


# Example usage
ppt_path = "Chapter-1 (1).pptx"  # Replace with the path to your PowerPoint file

# Set USE_API to True to use the Hugging Face API or False for local model
summary = summarize_ppt(ppt_path, USE_API=True, Save_to_txt=True)  # Set USE_API to False for local summarization
print("Summary:")
print(summary)
